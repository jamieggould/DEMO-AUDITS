"""
Lawmens Pre-Demolition Audit Generator — Flask backend
Template: Savills-7.pptx
"""
import os
import io
import re
import base64
import traceback
from datetime import datetime
from flask import Flask, render_template, request, jsonify, send_file, Response
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
try:
    from PIL import Image as _PILImage
    _PIL_AVAILABLE = True
except ImportError:
    _PIL_AVAILABLE = False

load_dotenv()

app = Flask(__name__)
app.secret_key = os.environ.get("SECRET_KEY", "lawmens-audit-2024")

OPENAI_API_KEY     = os.environ.get("OPENAI_API_KEY")
PPTX_TEMPLATE_PATH = os.environ.get("PPTX_TEMPLATE_PATH", "Savills-7.pptx")
APP_PASSWORD       = os.environ.get("APP_PASSWORD", "Lawmens123")

# ─────────────────────────────────────────────────────────────
# HTTP BASIC AUTH — password-protect every route except /health
# ─────────────────────────────────────────────────────────────

@app.before_request
def _require_password():
    if request.path == '/health':
        return  # let uptime monitors through unauthenticated
    auth = request.headers.get('Authorization', '')
    if auth.startswith('Basic '):
        try:
            credentials = base64.b64decode(auth[6:]).decode('utf-8')
            _user, password = credentials.split(':', 1)
            if password == APP_PASSWORD:
                return  # ✓ authenticated
        except Exception:
            pass
    return Response(
        'Access restricted. Please enter your credentials.',
        401,
        {'WWW-Authenticate': 'Basic realm="Lawmens Audit Generator"'},
    )

# ─────────────────────────────────────────────────────────────
# DEFAULT TEXT & EWC CODES FOR EACH MATERIAL TYPE
# ─────────────────────────────────────────────────────────────
MATERIAL_DEFAULTS = {
    'Carpet': {
        'ecf':          3.40,
        'description':  (
            'Carpet tiles and broadloom carpet removed during strip-out operations. '
            'Condition varies by area of use and age, with tiles in lower-traffic zones '
            'typically retaining sufficient integrity for direct reuse. All carpet is to be '
            'inspected prior to removal to assess reuse viability.'
        ),
        'waste_rec':    (
            'Carpet tiles in good, undamaged condition to be carefully uplifted, palletised '
            'and offered for reuse via material exchange platforms such as Carpet Recycling UK '
            'or donated to local charities and community organisations. Worn, stained or '
            'heavily soiled carpet to be segregated and sent to a specialist textile recycler '
            'for fibre recovery or energy recovery as a final option.'
        ),
        'risks':        (
            'Adhesive residue on tile backing may reduce reuse appeal and complicate recycling. '
            'Market demand for second-hand carpet tiles is highly condition-dependent and may '
            'require advance advertising via platforms such as Salvo or Freecycle. Some older '
            'carpets may contain hazardous fibres or fire-retardant chemicals requiring specialist disposal.'
        ),
        'potential':    'Medium',
        'ewc':          '20 03 01',
    },
    'Timber': {
        'ecf':          0.46,
        'description':  (
            'Structural and non-structural timber elements including studwork, boarding, door '
            'frames, skirtings, and secondary timber framing. Timber quality varies across the '
            'building — elements in accessible, dry locations are most likely to be suitable '
            'for reuse. All timber is to be visually assessed for rot, infestation or damage '
            'prior to selection for salvage.'
        ),
        'waste_rec':    (
            'Good-quality structural timber to be carefully dismantled, de-nailed and offered '
            'for reuse on-site, on other construction projects, or sold through architectural '
            'salvage merchants. Timber that cannot be directly reused but is clean and untreated '
            'to be sent for biomass energy recovery. Painted, treated or contaminated timber '
            'to be segregated and sent to a licensed wood recycler with appropriate classification.'
        ),
        'risks':        (
            'Older timber may carry lead-based paint requiring specialist assessment before handling. '
            'Embedded fixings and nails must be removed prior to reuse, increasing processing cost. '
            'Moisture damage, dry rot or insect infestation may significantly reduce structural '
            'reuse potential and may require licensed disposal. Preservative-treated timber '
            'may carry hazardous substance restrictions.'
        ),
        'potential':    'High',
        'ewc':          '17 02 01',
    },
    'Plasterboard': {
        'ecf':          0.39,
        'description':  (
            'Gypsum plasterboard lining, partitions and ceiling boards including standard, '
            'moisture-resistant and fire-rated grades. Boards are typically fixed with screws '
            'and/or adhesive and will be in varying condition depending on installation method '
            'and exposure. Clean, dry boards represent the highest recycling value.'
        ),
        'waste_rec':    (
            'Clean, segregated plasterboard to be sent to a licensed plasterboard recycler for '
            'closed-loop gypsum recovery — operators such as British Gypsum and Siniat operate '
            'take-back schemes. All plasterboard must be kept dry, free from adhesive contamination '
            'and segregated from other waste streams to maximise recycling yield and minimise '
            'processing costs. Contaminated or mixed boards to be disposed of via licensed '
            'waste contractor as a last resort.'
        ),
        'risks':        (
            'Contamination with paint, adhesive, tile grout or other bonded materials will '
            'prevent acceptance at specialist recyclers. Boards that have been wet may contain '
            'elevated sulphate content, restricting disposal routes. Older plasterboard may '
            'contain asbestos in the skim coat — a survey is required before any works commence. '
            'Breakage during removal significantly increases processing costs.'
        ),
        'potential':    'Low',
        'ewc':          '17 08 02',
    },
    'Glass': {
        'ecf':          0.91,
        'description':  (
            'Glazing units, frameless partition screens, mirrors, curtain walling infill panels '
            'and internally glazed door panels removed during strip-out. Units vary from single '
            'glazed to double-glazed sealed units depending on building age and location. '
            'Condition is to be assessed carefully to determine reuse suitability.'
        ),
        'waste_rec':    (
            'Intact, undamaged glass panels and screens to be carefully removed with appropriate '
            'handling equipment and offered for reuse through architectural salvage or glass '
            'merchants. Double-glazed units with intact seals may have direct reuse value. '
            'Broken, laminated or coated glass to be segregated by type and sent to a '
            'specialist glass cullet recycler for reprocessing into new glass products.'
        ),
        'risks':        (
            'Significant safety hazard during removal — all personnel must wear appropriate '
            'cut-resistant PPE. Laminated, toughened or coated glass cannot be recycled through '
            'standard flat glass routes and requires specialist handling. Transport requires '
            'purpose-built A-frames and secure packaging. Some older glass may contain lead '
            'used in seals or coatings.'
        ),
        'potential':    'Medium',
        'ewc':          '17 02 02',
    },
    'Metal': {
        'ecf':          1.37,
        'description':  (
            'Structural and non-structural metal including raised access floor systems, '
            'suspended ceiling grid, metal stud framing, pipework, ductwork, cable trays, '
            'and miscellaneous ferrous and non-ferrous fixings and fittings. Metal represents '
            'one of the highest-value recyclable waste streams on the project.'
        ),
        'waste_rec':    (
            'All scrap metal to be carefully segregated by type (ferrous/non-ferrous) on-site '
            'and collected by a licensed metal recycler. Segregated streams attract significantly '
            'higher recycling value than mixed metal. Raised access floor panels in good, '
            'undamaged condition to be offered for direct reuse through flooring salvage '
            'specialists. Structural steel elements to be assessed for direct reuse potential '
            'before recycling.'
        ),
        'risks':        (
            'Mixed metal streams substantially reduce recycling income. Sharp edges, cut ends '
            'and residual fixings create handling hazards requiring appropriate PPE throughout. '
            'Some older components may have asbestos-based insulation or fire-retardant coatings — '
            'an asbestos survey is required before any works begin. Refrigerant-containing plant '
            'items must be handled by an F-Gas registered contractor.'
        ),
        'potential':    'High',
        'ewc':          '17 04 05',
    },
    'Hardcore': {
        'ecf':          0.13,
        'description':  (
            'Inert masonry, concrete, floor screeds, ceramic and porcelain tiles, and brick '
            'materials arising from internal demolition and strip-out works. Volumes will be '
            'dependent on the extent of structural interventions. Clean, segregated hardcore '
            'has strong recycling demand as secondary aggregate.'
        ),
        'waste_rec':    (
            'Clean hardcore materials to be broken down, segregated and crushed on-site or '
            'sent to a licensed inert waste recycler for reprocessing as secondary aggregate '
            'for use in sub-base or concrete production. Ceramic tiles in good condition to '
            'be offered for direct reuse. Brick to be cleaned of mortar and assessed for '
            'salvage value prior to disposal. All inert waste must be transferred under a '
            'valid waste transfer note to an appropriately permitted facility.'
        ),
        'risks':        (
            'Contamination with adhesive, render, paint or other bonded materials reduces '
            'recyclability and may reclassify waste as non-inert, triggering higher disposal '
            'costs. Composite tiles with vinyl or rubber backings require separation before '
            'recycling. Older grout or tile adhesive may contain hazardous substances.'
        ),
        'potential':    'Low',
        'ewc':          '17 01 01',
    },
    'Insulation': {
        'ecf':          1.28,
        'description':  (
            'Thermal and acoustic insulation materials including mineral wool batts and slabs, '
            'rigid foam boards (EPS, XPS, PIR/PUR), pipe insulation lagging and cavity fill. '
            'Condition and reuse suitability will vary with installation method and exposure '
            'to moisture or mechanical damage during the building\'s operational life.'
        ),
        'waste_rec':    (
            'Clean, undamaged rigid insulation boards in good condition to be carefully removed '
            'and offered for reuse on other construction or refurbishment projects via surplus '
            'material trading platforms. Mineral wool to be bagged and sent to a specialist '
            'mineral wool recycler. Manufacturer take-back schemes (e.g. Rockwool, Isover, '
            'Kingspan) should be explored. Contaminated or wet insulation to be disposed of '
            'via licensed contractor.'
        ),
        'risks':        (
            'Fibrous mineral wool and refractory ceramic fibre (RCF) products require appropriate '
            'respiratory PPE during handling and may be classified as hazardous waste if RCF '
            'content is confirmed. Some older spray foam insulation may contain ozone-depleting '
            'blowing agents — specialist contractor and disposal is required. Rigid foam boards '
            'with fire-retardant coatings may have restricted disposal routes.'
        ),
        'potential':    'Low',
        'ewc':          '17 06 04',
    },
    'Fibre Ceiling Tiles': {
        'ecf':          0.94,
        'description':  (
            'Suspended mineral fibre, acoustic and glass fibre ceiling tiles from suspended '
            'grid systems, typically 600x600mm and 600x1200mm format. Tiles are acoustically '
            'and thermally functional but fragile — condition is highly dependent on age, '
            'moisture exposure and handling during previous installation and removal works.'
        ),
        'waste_rec':    (
            'Undamaged, clean tiles to be carefully removed from the grid, stacked back-to-back '
            'and palletised for reuse or donation to community or refurbishment projects. '
            'Damaged, stained or watermarked tiles to be sent to a specialist ceiling tile '
            'recycler such as Armstrong Ceiling Solutions or Saint-Gobain Ecophon, who operate '
            'closed-loop recycling schemes. Tiles must be kept dry and free of contamination.'
        ),
        'risks':        (
            'Fragile nature results in a typical breakage rate of 15–30% during removal, '
            'significantly reducing the reusable yield. Moisture staining, mould growth or '
            'watermarks will prevent reuse. Some older mineral fibre tiles may contain asbestos — '
            'a comprehensive asbestos survey must be completed before removal commences. '
            'Tiles with integral lighting or service penetrations may not be accepted by recyclers.'
        ),
        'potential':    'Medium',
        'ewc':          '17 06 05',
    },
    'Plastic': {
        'ecf':          2.00,
        'description':  (
            'Mixed plastic components including conduit, cable management trunking, switch '
            'and socket faceplates, signage, pipe fittings, and miscellaneous plastic fittings '
            'and finishes. Plastic waste streams are diverse in polymer type, which significantly '
            'affects recyclability and market value.'
        ),
        'waste_rec':    (
            'Segregate plastic waste by polymer type where practicable (ABS, PVC, PP, HDPE) '
            'to maximise recycling value and enable specialist processing. Large, clean plastic '
            'items in good condition to be offered for direct reuse. PVC cable management and '
            'conduit to be sent to a licensed PVC recycler. Mixed, unidentified plastic to be '
            'segregated from other waste streams and sent to a licensed plastics recycler for '
            'sorting and reprocessing.'
        ),
        'risks':        (
            'Mixed, unsegregated plastic streams attract very low or zero recycling value. '
            'PVC materials may contain plasticisers and stabilisers that require specialist '
            'handling and controlled processing. Some older plastics may contain hazardous '
            'additives (e.g. cadmium, lead-based stabilisers in older PVC). Certain reinforced '
            'or composite plastics cannot be recycled through standard routes.'
        ),
        'potential':    'Low',
        'ewc':          '17 02 03',
    },
    'Vinyl': {
        'ecf':          2.20,
        'description':  (
            'Vinyl sheet flooring and luxury vinyl tiles (LVT) removed during strip-out, '
            'including both loose-lay and fully adhered installations. Sheet vinyl is typically '
            '2–3.5mm thick with a PVC wear layer and fibreglass or felt backing. Condition '
            'varies with age, adhesive type and traffic levels.'
        ),
        'waste_rec':    (
            'Vinyl flooring in good condition to be rolled and offered for reuse on other '
            'projects or donated through material exchange networks. All vinyl flooring to be '
            'referred to manufacturer take-back or closed-loop recycling schemes — Altro, '
            'Tarkett and Polyflor operate vinyl recycling programmes. Worn, torn or heavily '
            'contaminated vinyl to be sent to a licensed vinyl recycler for PVC recovery.'
        ),
        'risks':        (
            'Adhesive backing compounds in older installations may contain hazardous substances '
            'including solvents and bituminous materials, restricting disposal and recycling routes. '
            'Floor tiles installed prior to 2000 may contain asbestos — a comprehensive survey '
            'must be completed before any works commence. Some vinyl-backed tiles cannot be '
            'separated from the substrate without contamination.'
        ),
        'potential':    'Medium',
        'ewc':          '20 01 39',
    },
    'Rubber': {
        'ecf':          2.10,
        'description':  (
            'Rubber flooring, anti-vibration isolation mounts, door seals, gaskets and '
            'acoustic isolation pads removed during strip-out. Rubber flooring is typically '
            'found in areas requiring slip resistance or acoustic control. Material may be '
            'bonded, loose-lay or mechanically fixed.'
        ),
        'waste_rec':    (
            'Rubber flooring in good, structurally sound condition to be offered for reuse in '
            'sports, educational or industrial applications. Manufacturers such as Altro and '
            'Nora operate rubber take-back schemes. Scrap rubber to be collected by a '
            'specialist rubber recycler for crumb rubber production for use in playgrounds, '
            'sports surfaces and road surfaces. Anti-vibration mounts and gaskets to be '
            'assessed individually for reuse potential.'
        ),
        'risks':        (
            'Adhesive-bonded rubber installations are extremely difficult to remove cleanly '
            'without substrate contamination, significantly reducing reuse potential. '
            'Synthetic rubber compounds may contain carbon black, plasticisers and other '
            'additives that restrict disposal routes. Some rubber products may contain '
            'polychlorinated biphenyls (PCBs) in older adhesive formulations.'
        ),
        'potential':    'Medium',
        'ewc':          '16 01 03',
    },
    'Fabric': {
        'ecf':          5.50,
        'description':  (
            'Textile materials including window blinds and tracks, curtains and rails, '
            'upholstered furniture fabric, acoustic wall panels and pin boards. Fabric waste '
            'is typically heterogeneous in fibre composition, which affects the available '
            'recycling routes and market demand.'
        ),
        'waste_rec':    (
            'Clean, undamaged fabric items and soft furnishings in good condition to be '
            'donated to registered charities, community organisations or social enterprises. '
            'Window blinds in working condition to be offered on material exchange platforms. '
            'Contaminated, worn or mixed-fibre fabrics to be sent to a textile recycler for '
            'fibre recovery or industrial cleaning cloth manufacture. Acoustic panels to be '
            'assessed for manufacturer take-back schemes.'
        ),
        'risks':        (
            'Fire-retardant treatments applied to many commercial textiles may contain PFAS '
            'or other hazardous chemicals, restricting disposal and reuse routes. Heavily '
            'soiled, stained or odour-contaminated items are unlikely to find reuse markets. '
            'Mixed-fibre compositions (e.g. polyester-cotton blends) significantly reduce '
            'recycling options compared to single-fibre streams.'
        ),
        'potential':    'Low',
        'ewc':          '20 01 10',
    },
    'Fluorescent Tubes': {
        'ecf':          1.50,
        'description':  (
            'Linear fluorescent lamps (T5, T8, T12 formats), compact fluorescent lamps (CFLs) '
            'and associated control gear removed during de-fit. Fluorescent lamps contain '
            'mercury and are classified as hazardous waste, requiring specialist handling, '
            'storage and disposal via a licensed WEEE contractor throughout all stages of '
            'the project.'
        ),
        'waste_rec':    (
            'All fluorescent and compact fluorescent lamps must be collected, handled and '
            'stored in approved lamp containers by a licensed waste electrical and electronic '
            'equipment (WEEE) contractor for specialist mercury recovery and glass recycling. '
            'Lamps must not be crushed, broken or mixed with general waste at any stage. '
            'A licensed waste transfer note referencing the relevant EWC code must be obtained '
            'for all lamp disposals. Lamp disposal must be included in the site waste management plan.'
        ),
        'risks':        (
            'Fluorescent tubes contain elemental mercury — a highly hazardous substance. '
            'Breakage creates an acute mercury vapour inhalation hazard requiring immediate '
            'evacuation and specialist clean-up. Lamps must not be disposed of in general '
            'skip waste under any circumstances. Older T12 tubes may contain higher mercury '
            'concentrations than modern equivalents. Hazardous waste consignment notes and '
            'licensed carrier documentation are mandatory.'
        ),
        'potential':    'Low',
        'ewc':          '20 01 21',
    },
    'Oil / Hydraulic Fluid': {
        'ecf':          0.45,
        'description':  (
            'Hydraulic oils and lubricants from mechanical plant, lift hydraulic systems, '
            'compressors, generators and building services equipment. Waste oil is classified '
            'as hazardous waste and requires dedicated containment, manifesting and licensed '
            'contractor collection at all stages from the point of generation through to '
            'final disposal or re-refining.'
        ),
        'waste_rec':    (
            'All waste oils to be drained and collected in sealed, dedicated containers by a '
            'licensed waste oil contractor for re-refining into recycled base oil or '
            'energy recovery as a secondary option. A hazardous waste consignment note must '
            'be completed for every collection. Waste oil must be stored in a bunded '
            'secondary containment area prior to collection, with a minimum containment '
            'capacity of 110% of the largest container volume. No mixing with other '
            'waste streams is permitted.'
        ),
        'risks':        (
            'Classified as hazardous waste under the Hazardous Waste Regulations — '
            'a licensed waste carrier and appropriately permitted disposal facility '
            'are mandatory. Significant spillage risk during draining and storage. '
            'Contamination of surrounding materials with waste oil will result in '
            'reclassification of those materials as hazardous. Hydraulic lift systems '
            'may contain chlorinated solvents or other environmentally harmful additives '
            'requiring specialist identification before disposal.'
        ),
        'potential':    'Low',
        'ewc':          '13 01 10',
    },
}

EWC_CODES = {k: v['ewc'] for k, v in MATERIAL_DEFAULTS.items()}

# ─────────────────────────────────────────────────────────────
# BASIC ROUTES
# ─────────────────────────────────────────────────────────────

@app.route("/")
def index():
    return render_template("index.html")

@app.route("/health")
def health():
    return "OK", 200

# ─────────────────────────────────────────────────────────────
# OPENAI TEXT GENERATION
# ─────────────────────────────────────────────────────────────

def openai_generate(prompt):
    """
    Generate text via OpenAI.
    Returns (text, None) on success, or (None, error_string) on failure.
    """
    if not OPENAI_API_KEY:
        return None, "OPENAI_API_KEY is not set in the environment."
    try:
        from openai import OpenAI
        client = OpenAI(api_key=OPENAI_API_KEY)
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}]
        )
        return response.choices[0].message.content.strip(), None
    except ImportError:
        return None, "The 'openai' Python package is not installed. Add it to requirements.txt."
    except Exception as e:
        traceback.print_exc()
        return None, str(e)

@app.route("/generate-ai-text", methods=["POST"])
def generate_ai_text():
    data    = request.get_json()
    section = data.get("section")
    report  = data.get("report_data", {})
    addr    = report.get("job_address", "the site")
    client  = report.get("client_name", "the client")
    mats    = report.get("kwp_materials", [])
    mat_str = ", ".join(m.get("name", "") for m in mats if m.get("name"))

    prompts = {
        "executive_summary": (
            f"Write a detailed, professional executive summary (4–5 paragraphs, approximately 350–450 words) "
            f"for a pre-refurbishment audit report prepared by Lawmens at {addr} for client {client}. "
            f"Key waste products identified during the site survey include: {mat_str}. "
            f"The summary should: (1) introduce the purpose and scope of the pre-refurbishment audit; "
            f"(2) describe the site visit process and how materials were assessed; "
            f"(3) summarise the key findings including the types and approximate quantities of key waste products; "
            f"(4) outline the circular economy principles being applied, including reuse, recycling and "
            f"landfill diversion targets; (5) conclude with Lawmens' overall recommendation to set a "
            f"minimum 95% landfill diversion target and maximise reuse opportunities. "
            f"Use formal, professional language appropriate for a sustainability report submitted "
            f"to a planning authority or client team. Do not use bullet points."
        ),
        "conclusion": (
            f"Write a detailed, professional conclusion and recommendations section (3–4 paragraphs, "
            f"approximately 250–350 words) for a pre-refurbishment audit at {addr} prepared by Lawmens. "
            f"Materials identified: {mat_str}. "
            f"The section should: (1) summarise the overall findings of the audit and the total "
            f"estimated material volumes; (2) set out specific, measurable recommendations including "
            f"a minimum 95% diversion from landfill, minimum 20% by volume or 30% by weight for "
            f"direct reuse, and a target recycling rate of at least 80%; "
            f"(3) recommend that a Site Waste Management Plan (SWMP) be prepared and implemented "
            f"by the principal contractor prior to works commencing; "
            f"(4) state that actual material arisings should be compared against forecast figures "
            f"during the strip-out phase to monitor performance. "
            f"Use formal, professional language. Do not use bullet points."
        ),
        "circular_economy": (
            f"Write 3–4 sentences describing exemplary circular economy commitments for a "
            f"construction refurbishment project at {addr}. The text should reference: "
            f"material reuse targets (minimum 20% by volume), recycling targets (minimum 80%), "
            f"landfill diversion (minimum 95%), GLA Whole Life Carbon benchmarks, "
            f"circular economy principles aligned with the London Plan, "
            f"and the use of a Site Waste Management Plan. "
            f"Write in the first person plural (we/our) as if written by the project team. "
            f"Keep it concise and suitable for insertion mid-paragraph in a technical report."
        ),
        "information_provided": (
            f"Write a concise but detailed description (2–3 sentences followed by a comma-separated list) "
            f"of the types of information and documents typically provided to the auditor prior to "
            f"and during a pre-refurbishment audit at {addr}. "
            f"Include: architectural drawings (as-built where available), structural drawings, "
            f"M&E services drawings, the existing asbestos management survey or refurbishment and "
            f"demolition (R&D) asbestos survey, floor area schedules, previous occupier information, "
            f"photographic records from the site visit, and any relevant planning documents. "
            f"Write in a formal, professional tone suitable for a planning report."
        ),
        "benchmark": (
            f"Write one concise but informative sentence describing the resource efficiency benchmark "
            f"being adopted for the project at {addr}, referencing the GLA Whole Life Carbon (WLC) "
            f"benchmarks, the WRAP vision for a circular economy in construction, and the target of "
            f"diverting at least 95% of demolition and construction waste from landfill in line with "
            f"best practice guidance. Keep it suitable for insertion within a technical paragraph."
        ),
        "aims_resource_efficiency": (
            f"Write a concise phrase (1–2 sentences, no more than 40 words) describing the resource "
            f"efficiency aims for the Site Waste Management Plan at {addr}. "
            f"Reference diverting at least 95% of waste from landfill, maximising material reuse on "
            f"and off site, and tracking actual versus forecast material arisings. "
            f"The text will be inserted mid-sentence in a technical report, so write it as a "
            f"continuation — start from 'diverting' or a similar neutral entry point."
        ),
    }

    if section not in prompts:
        return jsonify({"error": "Unknown section"}), 400

    text, err = openai_generate(prompts[section])
    if err:
        return jsonify({"error": err}), 500
    return jsonify({"text": text})

@app.route("/generate-material-text", methods=["POST"])
def generate_material_text():
    """Generate description/recommendation/risks for a single material."""
    data     = request.get_json()
    mat_name = data.get("material_name", "")
    field    = data.get("field", "description")

    # Use defaults if available
    if mat_name in MATERIAL_DEFAULTS and field in MATERIAL_DEFAULTS[mat_name]:
        return jsonify({"text": MATERIAL_DEFAULTS[mat_name][field]})

    # Fall back to AI
    prompts = {
        "description": (
            f"Write 3–4 detailed sentences describing {mat_name} as a waste material in a "
            f"pre-refurbishment audit context. Include typical quantities, condition assessment "
            f"approach, and what factors affect its reuse or recycling potential."
        ),
        "waste_rec": (
            f"Write 3–4 detailed sentences with waste management recommendations for {mat_name} "
            f"in a commercial demolition or strip-out project. Include preferred reuse routes, "
            f"specific recycling options and relevant industry schemes or contractors. "
            f"Reference the waste hierarchy."
        ),
        "risks": (
            f"Write 3–4 detailed sentences describing the key risk factors and constraints for "
            f"reusing or recycling {mat_name} from a commercial refurbishment project. Include "
            f"hazardous substance risks, logistical challenges, and market or quality limitations."
        ),
    }
    prompt = prompts.get(field, prompts["description"])
    text, err = openai_generate(prompt)
    if err:
        return jsonify({"error": err}), 500
    return jsonify({"text": text})

# ─────────────────────────────────────────────────────────────
# EWC LOOKUP
# ─────────────────────────────────────────────────────────────

@app.route("/ewc-lookup", methods=["POST"])
def ewc_lookup():
    name     = request.get_json().get("name", "")
    ewc      = EWC_CODES.get(name, "")
    defaults = MATERIAL_DEFAULTS.get(name, {})
    return jsonify({
        "ewc":         ewc,
        "ecf":         defaults.get("ecf", ""),
        "description": defaults.get("description", ""),
        "waste_rec":   defaults.get("waste_rec", ""),
        "risks":       defaults.get("risks", ""),
        "potential":   defaults.get("potential", "Medium"),
    })

# ─────────────────────────────────────────────────────────────
# EXCEL CALCULATOR PARSER
# ─────────────────────────────────────────────────────────────

CALC_MATERIAL_MAP = [
    (20, 'Carpet',            'Timber'),
    (26, 'Plasterboard',      'Glass'),
    (32, 'Metal',             'Hardcore'),
    (38, 'Insulation',        'Fibre Ceiling Tiles'),
    (44, 'Plastic',           'Vinyl'),
    (50, 'Rubber',            'Fabric'),
    (56, 'Fluorescent Tubes', 'Oil / Hydraulic Fluid'),
]

def parse_calculator_excel(file_bytes):
    from openpyxl import load_workbook
    wb   = load_workbook(io.BytesIO(file_bytes), data_only=True)
    ws   = wb['Sheet1']
    rows = list(ws.iter_rows(values_only=True))

    materials = []
    for base_row, left_name, right_name in CALC_MATERIAL_MAP:
        val_row = rows[base_row + 2]
        pct_row = rows[base_row + 4]
        for name, wt_col, vol_col, pct_col in [
            (left_name,  3, 5, 3),
            (right_name, 9, 11, 9),
        ]:
            materials.append({
                'name':       name,
                'weight_kg':  float(val_row[wt_col]  or 0),
                'volume_m3':  float(val_row[vol_col] or 0),
                'weight_pct': round(float(pct_row[pct_col] or 0) * 100, 1),
                'ewc':        EWC_CODES.get(name, ''),
                'weight_t':   round(float(val_row[wt_col] or 0) / 1000, 3),
            })

    total_row = rows[66]
    return {
        'materials':       materials,
        'total_weight_t':  round(float(total_row[2] or 0), 3),
        'total_volume_m3': round(float(total_row[8] or 0), 3),
    }

@app.route("/parse-calculator", methods=["POST"])
def parse_calculator():
    f = request.files.get("calculator_file")
    if not f:
        return jsonify({"error": "No file uploaded"}), 400
    try:
        return jsonify(parse_calculator_excel(f.read()))
    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

# ─────────────────────────────────────────────────────────────
# PPTX — TEXT REPLACEMENT (handles text boxes AND table cells)
# ─────────────────────────────────────────────────────────────

def _sanitise(val):
    """
    Convert a replacement value to a clean string for PPTX XML.
    Strips carriage-return characters that PPTX encodes as _x000D_ in the output.
    """
    s = str(val) if val is not None else ''
    return s.replace('\r\n', '\n').replace('\r', '')

def _replace_in_paragraph(para, replacements):
    """
    Merge all runs in a paragraph, replace {{KEY}} tokens, write back.
    Handles the common PPTX split-run problem.
    """
    if not para.runs:
        return
    full = ''.join(r.text for r in para.runs)
    new  = full
    for key, val in replacements.items():
        new = new.replace(f'{{{{{key}}}}}', _sanitise(val))
    if new != full:
        para.runs[0].text = new
        for r in para.runs[1:]:
            r.text = ''

def _replace_in_text_frame(tf, replacements):
    for para in tf.paragraphs:
        _replace_in_paragraph(para, replacements)

def _replace_in_shape(shape, replacements):
    """Replace in text frames, table cells, and grouped shapes."""
    if shape.has_text_frame:
        _replace_in_text_frame(shape.text_frame, replacements)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                _replace_in_text_frame(cell.text_frame, replacements)
    if shape.shape_type == 6:  # GROUP
        for s in shape.shapes:
            _replace_in_shape(s, replacements)

# ─────────────────────────────────────────────────────────────
# PPTX — IMAGE CONVERSION
# ─────────────────────────────────────────────────────────────

def _normalise_image(img_bytes):
    """
    Convert any uploaded file to PNG bytes for python-pptx.
    Handles: PDF (first page rendered via PyMuPDF), HEIC, WEBP, TIFF, BMP,
    and any other format Pillow can read. Falls back to the original bytes
    if all conversions fail.
    """
    # ── PDF: render first page to PNG via PyMuPDF ─────────────
    if img_bytes[:4] == b'%PDF':
        try:
            import fitz  # PyMuPDF
            doc  = fitz.open(stream=img_bytes, filetype='pdf')
            page = doc[0]
            # 2× zoom gives ~150 dpi — good quality without huge file size
            pix  = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            return pix.tobytes('png')
        except Exception:
            pass  # fall through to Pillow attempt

    # ── All other formats: normalise to PNG via Pillow ─────────
    if not _PIL_AVAILABLE:
        return img_bytes
    try:
        with _PILImage.open(io.BytesIO(img_bytes)) as img:
            if img.mode not in ('RGB', 'RGBA', 'L'):
                img = img.convert('RGB')
            out = io.BytesIO()
            img.save(out, format='PNG')
            out.seek(0)
            return out.read()
    except Exception:
        return img_bytes

# ─────────────────────────────────────────────────────────────
# PPTX — IMAGE REPLACEMENT
# ─────────────────────────────────────────────────────────────

def _replace_image_placeholders(prs, image_data):
    """
    image_data = { 'KEY': bytes }               — single image
               = { 'KEY': [bytes, bytes, ...] } — multiple (e.g. MAT_N_PHOTOS)
    When a key appears multiple times on a slide, each occurrence
    gets the next photo in the list.
    """
    img_map = {}
    for key, val in image_data.items():
        img_map[key] = [val] if isinstance(val, (bytes, bytearray)) else list(val)

    for slide in prs.slides:
        to_remove = []
        to_add    = []
        key_usage = {}

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            for key, photos in img_map.items():
                if f'{{{{{key}}}}}' in text and photos:
                    usage = key_usage.get(key, 0)
                    idx   = min(usage, len(photos) - 1)
                    to_remove.append(shape)
                    to_add.append((shape.left, shape.top, shape.width, shape.height, photos[idx]))
                    key_usage[key] = usage + 1
                    break

        for shape in to_remove:
            shape._element.getparent().remove(shape._element)
        for left, top, w, h, img_bytes in to_add:
            slide.shapes.add_picture(io.BytesIO(_normalise_image(img_bytes)), left, top, w, h)

# ─────────────────────────────────────────────────────────────
# PPTX — KWP DONUT CHART REPLACEMENT
# Outside labels "Name (X.X%)" — small slices (<4.5%) unlabelled
# ─────────────────────────────────────────────────────────────

_LABEL_MIN_PCT = 2.5   # only skip slices too tiny to label cleanly

def _configure_donut_labels(chart, pcts):
    """
    Write the <c:dLbls> XML block directly onto the plot element so that:
      - labels show the category name (two-line: name then percentage)
      - labels are positioned outside the slice with leader lines
      - slices below _LABEL_MIN_PCT have their label deleted
      - font is 8pt Calibri, same dark body-text colour, no bold
    """
    from lxml import etree
    from pptx.oxml.ns import qn

    plot_el = chart.plots[0]._element

    # Remove any existing dLbls so we start clean
    for old in plot_el.findall(qn('c:dLbls')):
        plot_el.remove(old)

    # OOXML schema order for doughnutChart:
    #   varyColors?, ser*, dLbls?, firstSliceAng?, holeSize?, extLst?
    # dLbls MUST be inserted before firstSliceAng/holeSize — appending at
    # the end breaks schema order and silently corrupts the chart.
    children  = list(plot_el)
    insert_at = len(children)  # fallback: end
    for i, child in enumerate(children):
        if child.tag in (qn('c:firstSliceAng'), qn('c:holeSize'), qn('c:extLst')):
            insert_at = i
            break

    dLbls = etree.Element(qn('c:dLbls'))
    plot_el.insert(insert_at, dLbls)

    # Delete individual labels for slices below the threshold
    for i, pct in enumerate(pcts):
        if pct < _LABEL_MIN_PCT:
            dLbl = etree.SubElement(dLbls, qn('c:dLbl'))
            etree.SubElement(dLbl, qn('c:idx')).set('val', str(i))
            etree.SubElement(dLbl, qn('c:delete')).set('val', '1')

    # Position: outside end (triggers leader lines automatically)
    etree.SubElement(dLbls, qn('c:dLblPos')).set('val', 'outEnd')

    # Show category name only (two-line value is encoded in the category string)
    for tag, val in [
        ('c:showLegendKey', '0'),
        ('c:showVal',       '0'),
        ('c:showCatName',   '1'),
        ('c:showSerName',   '0'),
        ('c:showPercent',   '0'),
    ]:
        etree.SubElement(dLbls, qn(tag)).set('val', val)

    # Font: 8pt Calibri, dark body-text colour (#404040), not bold, word-wrap on
    txPr   = etree.SubElement(dLbls, qn('c:txPr'))
    bodyPr = etree.SubElement(txPr, qn('a:bodyPr'))
    bodyPr.set('wrap', 'square')   # allow label text to wrap
    etree.SubElement(txPr, qn('a:lstStyle'))
    p    = etree.SubElement(txPr, qn('a:p'))
    pPr  = etree.SubElement(p,    qn('a:pPr'))
    defR = etree.SubElement(pPr,  qn('a:defRPr'))
    defR.set('sz', '800')          # 8pt — readable but compact
    defR.set('b',  '0')
    sf = etree.SubElement(defR, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr')).set('val', '404040')
    etree.SubElement(defR, qn('a:latin')).set('typeface', 'Calibri')


def _add_kwp_pie_chart(slide, left, top, width, height, mats, value_key):
    """
    Donut chart: two-line outside labels (name / percentage) per slice.
    Slices below _LABEL_MIN_PCT get no label to keep the chart clean.
    No legend — labels carry all needed information.
    """
    visible = [m for m in mats if float(m.get(value_key) or 0) > 0] or mats
    values  = [float(m.get(value_key) or 0) for m in visible]
    total   = sum(values) or 1
    pcts    = [v / total * 100 for v in values]

    # Two-line label: name on first line, percentage on second
    # Tiny slices get an empty string so no label box appears
    labels = [
        f"{m.get('name', '')}\n{pct:.1f}%" if pct >= _LABEL_MIN_PCT else ''
        for m, pct in zip(visible, pcts)
    ]

    cd = ChartData()
    cd.categories = labels
    cd.add_series('', values)

    gf    = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, left, top, width, height, cd)
    chart = gf.chart

    # Legend on the right — colour swatches + "Name\nX.X%" entries
    try:
        chart.has_legend = True
        chart.legend.include_in_layout = False
        from pptx.enum.chart import XL_LEGEND_POSITION
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.font.size = Pt(8)
        chart.legend.font.bold = False
        chart.legend.font.name = 'Calibri'
    except Exception:
        pass

    # Outside slice labels (two-line: name / percentage)
    _configure_donut_labels(chart, pcts)

def _replace_kwp_chart_placeholders(prs, kwp_materials):
    if not kwp_materials:
        return
    mats = [m for m in kwp_materials if m.get('name')]
    if not mats:
        return

    CHART_MAP = {
        'KWP_OF_TOTAL_WEIGHT': 'weight_pct',
        'KWP_BY_VOL':          'volume_m3',
        'KWP_BY_TON':          'weight_t',
    }
    for slide in prs.slides:
        to_remove, to_add = [], []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            for key, val_col in CHART_MAP.items():
                if f'{{{{{key}}}}}' in text:
                    to_remove.append(shape)
                    to_add.append((shape.left, shape.top, shape.width, shape.height, val_col))
                    break
        for shape in to_remove:
            shape._element.getparent().remove(shape._element)
        for left, top, w, h, val_col in to_add:
            _add_kwp_pie_chart(slide, left, top, w, h, mats, val_col)

# ─────────────────────────────────────────────────────────────
# PPTX — TRIM UNUSED MATERIAL ROWS & SLIDES
# Must run BEFORE text replacement so {{MATERIAL_N}} tags are intact.
# ─────────────────────────────────────────────────────────────

# Matches {{MATERIAL_EWC_N}}, {{MATERIAL_N}}, {{MATERIAL_N_xxx}}, {{MAT_N_xxx}}
_MAT_PLACEHOLDER_RE = re.compile(
    r'\{\{(?:MATERIAL_EWC_(\d+)|MATERIAL_(\d+)|MAT_(\d+)_)'
)

def _row_full_text(row):
    parts = []
    for cell in row.cells:
        for para in cell.text_frame.paragraphs:
            parts.append(''.join(run.text for run in para.runs))
    return ''.join(parts)

def _slide_full_text(slide):
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                parts.append(''.join(run.text for run in para.runs))
        if shape.has_table:
            for row in shape.table.rows:
                parts.append(_row_full_text(row))
        if shape.shape_type == 6:
            for s in shape.shapes:
                if s.has_text_frame:
                    for para in s.text_frame.paragraphs:
                        parts.append(''.join(run.text for run in para.runs))
    return ''.join(parts)

def _mat_indices_in_text(text):
    indices = set()
    for m in _MAT_PLACEHOLDER_RE.finditer(text):
        idx_str = next((g for g in m.groups() if g), None)
        if idx_str:
            indices.add(int(idx_str))
    return indices

def _trim_material_table_rows(prs, mat_count):
    """Delete table rows referencing a material index > mat_count."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            table = shape.table
            rows_to_remove = []
            for row_idx, row in enumerate(table.rows):
                indices = _mat_indices_in_text(_row_full_text(row))
                if indices and all(idx > mat_count for idx in indices):
                    rows_to_remove.append(row_idx)
            for row_idx in sorted(set(rows_to_remove), reverse=True):
                tr = table.rows[row_idx]._tr
                tr.getparent().remove(tr)

def _trim_empty_material_slides(prs, mat_count):
    """Delete slides whose only material references are for index > mat_count."""
    from pptx.oxml.ns import qn as _qn
    slides_to_remove = []
    for slide_idx, slide in enumerate(prs.slides):
        text    = _slide_full_text(slide)
        indices = _mat_indices_in_text(text)
        if indices and all(idx > mat_count for idx in indices):
            slides_to_remove.append(slide_idx)

    xml_slides = prs.slides._sldIdLst
    for slide_idx in sorted(slides_to_remove, reverse=True):
        rId = xml_slides[slide_idx].get(_qn('r:id'))
        prs.part.drop_rel(rId)
        del xml_slides[slide_idx]

# Matches {{SPEC_1}} … {{SPEC_7}}
_SPEC_PLACEHOLDER_RE = re.compile(r'\{\{SPEC_(\d+)\}\}')

def _spec_indices_in_text(text):
    """Return set of SPEC slot numbers found in a text string."""
    return {int(m.group(1)) for m in _SPEC_PLACEHOLDER_RE.finditer(text)}

def _trim_unused_spec_slides(prs, provided_spec_indices):
    """
    Delete spec-image slides whose {{SPEC_N}} slot was not provided.
    provided_spec_indices: set of ints, e.g. {1, 2, 3} if 3 specs were uploaded.
    A slide is removed when ALL its SPEC_N references are outside the provided set.
    """
    from pptx.oxml.ns import qn as _qn
    slides_to_remove = []
    for slide_idx, slide in enumerate(prs.slides):
        text    = _slide_full_text(slide)
        indices = _spec_indices_in_text(text)
        # Only remove if the slide exists solely to show missing spec images
        if indices and all(idx not in provided_spec_indices for idx in indices):
            slides_to_remove.append(slide_idx)

    xml_slides = prs.slides._sldIdLst
    for slide_idx in sorted(slides_to_remove, reverse=True):
        rId = xml_slides[slide_idx].get(_qn('r:id'))
        prs.part.drop_rel(rId)
        del xml_slides[slide_idx]

# ─────────────────────────────────────────────────────────────
# FILL TEMPLATE
# ─────────────────────────────────────────────────────────────

def fill_pptx_template(replacements, image_data=None, kwp_materials=None, provided_spec_indices=None):
    prs       = Presentation(PPTX_TEMPLATE_PATH)
    mat_count = len(kwp_materials) if kwp_materials else 0
    # Trim unused rows/slides BEFORE replacement (placeholders must be intact)
    if mat_count < 30:
        _trim_material_table_rows(prs, mat_count)
        _trim_empty_material_slides(prs, mat_count)
    if provided_spec_indices is not None and len(provided_spec_indices) < 7:
        _trim_unused_spec_slides(prs, provided_spec_indices)
    for slide in prs.slides:
        for shape in slide.shapes:
            _replace_in_shape(shape, replacements)
    if image_data:
        _replace_image_placeholders(prs, image_data)
    if kwp_materials:
        _replace_kwp_chart_placeholders(prs, kwp_materials)
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out

# ─────────────────────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────────────────────

def _first(data, key, default=''):
    vals = data.get(key, [])
    if isinstance(vals, list):
        return vals[0] if vals else default
    return vals if vals else default

def _read_upload(fs):
    if not fs or fs.filename == '':
        return None
    return fs.read()

# ─────────────────────────────────────────────────────────────
# BUILD REPLACEMENTS
# Maps every {{PLACEHOLDER}} in Savills-7.pptx to a form value
# ─────────────────────────────────────────────────────────────

def build_replacements(data, mat_list):
    """
    mat_list: list of dicts (one per material, index 0 = material 1) with keys:
      name, ewc, vol, weigh, weighp, ecf, carb, reuse,
      waste_rec, desc, potential, risks
    """
    def g(k, d=''):
        return _first(data, k, d)

    r = {}

    # ── Cover / global ────────────────────────────────────────
    r['PROJECT_ADDRESS'] = g('job_address')
    r['CLIENT_NAME']     = g('client_name')
    r['DATE_OF_REPORT']  = g('date_of_report', datetime.now().strftime('%d %B %Y'))
    r['REPORT_NUMBER']   = g('report_number')

    # ── Team ──────────────────────────────────────────────────
    r['PREPARED_BY']        = g('prepared_by')
    r['PREPARED_BY_ROLE']   = g('prepared_by_role')
    r['PREPARED_DATE']      = g('prepared_date')
    r['AUTHORISED_BY']      = g('authorised_by')
    r['AUTHORISED_BY_ROLE'] = g('authorised_by_role')
    r['AUTHORISED_DATE']    = g('authorised_date')

    # ── Report narrative ──────────────────────────────────────
    r['CIRCULAR_ECONOMY_COMMITMENTS']      = g('circular_economy_commitments')
    r['BENCHMARK_FOR_RESOURCE_EFFICIENCY'] = g('benchmark_resource_efficiency')
    r['AIMS_RESOURCE_EFFICIENCY']          = g('aims_resource_efficiency')
    r['INFORMATION_PROVIDED']              = g('information_provided')
    r['KEY_WASTE_PRODUCTS']                = g('key_waste_products')
    r['PROJECT_WEIGHT']                    = g('project_weight')
    r['OVERALL_REUSE_PERCENT']             = g('overall_reuse_percent')
    r['LANDFILL_TARGET_PERCENT']           = g('landfill_target_percent', '95')
    r['RECYCLE_TARGET_PERCENT']            = g('recycle_target_percent',  '80')

    # ── Materials 1–30 ────────────────────────────────────────
    # Template uses three different EWC naming patterns:
    #   Slide 27 (rows 1-10):  MATERIAL_EWC_N
    #   Slides 28-29 (11-30):  MATERIAL_N_EWC
    #   Detail slides (31-45): MAT_N_EWC
    # We populate all three variants for every material.
    for i, m in enumerate(mat_list[:30], start=1):
        n = str(i)

        name    = m.get('name',     '')
        ewc     = m.get('ewc',      '')
        vol     = m.get('vol',      '')
        weigh   = m.get('weigh',    '')
        weighp  = m.get('weighp',   '')
        ecf     = m.get('ecf',      '')
        carb    = m.get('carb',     '')
        reuse   = m.get('reuse',    '')
        waste_r = m.get('waste_rec','')
        desc    = m.get('desc',     '')
        pot     = m.get('potential','Medium')
        risks   = m.get('risks',    '')

        # Material name (template also has 'MATERIAL 10' with space)
        r[f'MATERIAL_{n}'] = name
        if i == 10:
            r['MATERIAL 10'] = name

        # EWC — three naming patterns covering all slides
        r[f'MATERIAL_EWC_{n}'] = ewc   # slide 27 style (rows 1-10)
        r[f'MATERIAL_{n}_EWC'] = ewc   # slides 28-29 style (rows 11-30)
        r[f'MAT_{n}_EWC']      = ewc   # detail slide style

        # Numeric data
        r[f'MAT_{n}_VOL']    = vol
        r[f'MAT_{n}_WEIGH']  = weigh
        r[f'MAT_{n}_WEIGHP'] = weighp
        r[f'MAT_{n}_ECF']    = ecf
        r[f'MAT_{n}_CARB']   = carb
        r[f'MAT_{n}_REUSE']  = reuse

        # Text content
        r[f'MAT_{n}_WASTE_RECOMMENDATIONS'] = waste_r
        r[f'MATERIAL_{n}_DESCRIPTION']      = desc
        r[f'MATERIAL_{n}_POTENTIAL']        = pot
        r[f'MATERIAL_{n}_RISKS']            = risks

    # Blank out unfilled material slots so no {{PLACEHOLDER}} leaks through
    for i in range(len(mat_list) + 1, 31):
        n = str(i)
        for key in [
            f'MATERIAL_{n}',
            f'MATERIAL_EWC_{n}', f'MATERIAL_{n}_EWC', f'MAT_{n}_EWC',
            f'MAT_{n}_VOL', f'MAT_{n}_WEIGH', f'MAT_{n}_WEIGHP',
            f'MAT_{n}_ECF', f'MAT_{n}_CARB', f'MAT_{n}_REUSE',
            f'MAT_{n}_WASTE_RECOMMENDATIONS',
            f'MATERIAL_{n}_DESCRIPTION', f'MATERIAL_{n}_POTENTIAL', f'MATERIAL_{n}_RISKS',
        ]:
            r[key] = ''
        if i == 10:
            r['MATERIAL 10'] = ''

    return r

# ─────────────────────────────────────────────────────────────
# COLLECT IMAGES
# ─────────────────────────────────────────────────────────────

def _collect_image_data(files, mat_count):
    """
    Single images:   PREP_PHOTO, AUTH_PHOTO, BUILDING_PHOTO
    Spec images:     SPEC_1 … SPEC_7
    Material photos: MAT_N_PHOTOS → [photo1, photo2] per material
    """
    images = {}

    for field, key in [
        ('prepared_by_photo',   'PREP_PHOTO'),
        ('authorised_by_photo', 'AUTH_PHOTO'),
        ('building_photo',      'BUILDING_PHOTO'),
    ]:
        b = _read_upload(files.get(field))
        if b:
            images[key] = b

    for i in range(1, 8):
        b = _read_upload(files.get(f'spec_photo_{i}'))
        if b:
            images[f'SPEC_{i}'] = b

    for i in range(1, mat_count + 1):
        photos = []
        for j in (1, 2):
            b = _read_upload(files.get(f'mat_{i}_photo_{j}'))
            if b:
                photos.append(b)
        if photos:
            images[f'MAT_{i}_PHOTOS'] = photos

    return images

# ─────────────────────────────────────────────────────────────
# BUILD KWP MATERIAL LIST FROM FORM
# ─────────────────────────────────────────────────────────────

def _build_mat_list(data):
    """
    Read mat_N_* fields from form data into a list of dicts.
    Skips entries where name is blank.
    Auto-calculates carbon and weight% if missing.
    Supports up to 30 materials.
    """
    mat_list = []
    total_wt = 0.0

    # First pass — collect and compute total weight
    raw = []
    for i in range(1, 31):
        n    = str(i)
        name = _first(data, f'mat_{n}_name')
        if not name:
            continue
        wt = float(_first(data, f'mat_{n}_weigh', '0') or 0)
        total_wt += wt
        raw.append((i, name, wt))

    # Second pass — build full dicts with auto-calculated fields
    for i, name, wt in raw:
        n      = str(i)
        ecf    = float(_first(data, f'mat_{n}_ecf', '0') or 0)
        carb   = _first(data, f'mat_{n}_carb')
        weighp = _first(data, f'mat_{n}_weighp')

        if not carb and ecf and wt:
            carb = str(round(ecf * wt, 2))
        if not weighp and total_wt:
            weighp = str(round(wt / total_wt * 100, 1))

        defaults = MATERIAL_DEFAULTS.get(name, {})
        mat_list.append({
            'name':      name,
            'ewc':       _first(data, f'mat_{n}_ewc')       or defaults.get('ewc', ''),
            'vol':       _first(data, f'mat_{n}_vol',   ''),
            'weigh':     str(wt) if wt else '',
            'weighp':    weighp,
            'ecf':       _first(data, f'mat_{n}_ecf',   ''),
            'carb':      carb or '',
            'reuse':     _first(data, f'mat_{n}_reuse',  ''),
            'waste_rec': _first(data, f'mat_{n}_waste_rec') or defaults.get('waste_rec', ''),
            'desc':      _first(data, f'mat_{n}_desc')      or defaults.get('description', ''),
            'potential': _first(data, f'mat_{n}_potential') or defaults.get('potential', 'Medium'),
            'risks':     _first(data, f'mat_{n}_risks')     or defaults.get('risks', ''),
        })

    return mat_list, total_wt

# ─────────────────────────────────────────────────────────────
# MAIN REPORT GENERATION ROUTE
# ─────────────────────────────────────────────────────────────

@app.route("/generate-canva-report", methods=["POST"])
def generate_canva_report():
    data  = request.form.to_dict(flat=False)
    files = request.files

    mat_list, total_wt = _build_mat_list(data)

    # Auto-build KEY_WASTE_PRODUCTS if not manually set
    kwp_text = _first(data, 'key_waste_products')
    if not kwp_text and mat_list:
        kwp_text = ', '.join(m['name'] for m in mat_list)

    if kwp_text:
        data['key_waste_products'] = [kwp_text]
    if total_wt and not _first(data, 'project_weight'):
        data['project_weight'] = [f"{total_wt:.2f} tonnes"]

    replacements = build_replacements(data, mat_list)
    image_data   = _collect_image_data(files, len(mat_list))

    # Determine which spec slots have an image so unused spec slides can be dropped
    provided_spec_indices = {
        i for i in range(1, 8)
        if files.get(f'spec_photo_{i}') and files.get(f'spec_photo_{i}').filename
    }

    kwp_mats = [{
        'name':       m['name'],
        'weight_t':   float(m['weigh']  or 0),
        'weight_pct': float(m['weighp'] or 0),
        'volume_m3':  float(m['vol']    or 0),
    } for m in mat_list]

    try:
        output = fill_pptx_template(replacements, image_data, kwp_mats, provided_spec_indices)
    except FileNotFoundError:
        return jsonify({"error": (
            f"Template '{PPTX_TEMPLATE_PATH}' not found. "
            "Ensure Savills-7.pptx is committed to your repository root."
        )}), 500
    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": f"Failed to generate report: {str(e)}"}), 500

    addr     = _first(data, 'job_address').replace(' ', '_')[:40]
    filename = f"Audit_{addr}.pptx"
    return send_file(
        output,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name=filename,
    )

# ─────────────────────────────────────────────────────────────

if __name__ == "__main__":
    app.run(debug=True)
